
import csv
import os
import queue
import re
import subprocess
import sys
import threading
import traceback
import zipfile
from pathlib import Path
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from xml.etree import ElementTree as ET

# -----------------------------
# 선택 설치 라이브러리
# -----------------------------
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import xlrd
except ImportError:
    xlrd = None


APP_TITLE = "폴더 내부 키워드 검색기"
MAX_HITS_PER_FILE = 30

TEXT_EXTENSIONS = {
    ".txt", ".csv", ".tsv", ".log", ".ini", ".cfg", ".conf",
    ".json", ".xml", ".yaml", ".yml", ".md", ".py", ".java",
    ".c", ".cpp", ".h", ".hpp", ".js", ".ts", ".css", ".html",
    ".htm", ".sql", ".bat", ".cmd", ".ps1"
}

SUPPORTED_EXTENSIONS = (
    TEXT_EXTENSIONS
    | {".pdf", ".xlsx", ".xlsm", ".xltx", ".xltm", ".xls",
       ".docx", ".pptx", ".hwpx"}
)


def normalize_text(value):
    """검색 및 결과 표시용으로 문자열을 정리한다."""
    if value is None:
        return ""
    text = str(value)
    return re.sub(r"\s+", " ", text).strip()


def make_snippet(text, start, end, radius=70):
    """일치 위치 전후 문맥을 짧게 잘라 반환한다."""
    flat = normalize_text(text)
    if not flat:
        return ""
    start = max(0, min(start, len(flat)))
    end = max(start, min(end, len(flat)))
    left = max(0, start - radius)
    right = min(len(flat), end + radius)
    prefix = "…" if left > 0 else ""
    suffix = "…" if right < len(flat) else ""
    return prefix + flat[left:right] + suffix


def find_matches(text, pattern):
    """
    text 안에서 패턴을 찾아 (시작, 끝, 스니펫) 목록을 반환한다.
    pattern은 re.Pattern 객체다.
    """
    if text is None:
        return []
    text = str(text)
    results = []
    for match in pattern.finditer(text):
        results.append((match.start(), match.end(), make_snippet(text, match.start(), match.end())))
        if len(results) >= MAX_HITS_PER_FILE:
            break
    return results


def read_text_file(path):
    """텍스트 파일을 여러 인코딩으로 시도해 읽는다."""
    encodings = ("utf-8-sig", "utf-8", "cp949", "euc-kr", "utf-16")
    last_error = None

    for encoding in encodings:
        try:
            with open(path, "r", encoding=encoding, errors="strict") as f:
                return f.read()
        except UnicodeError as exc:
            last_error = exc

    # 마지막 수단: 읽을 수 없는 문자를 대체
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
    except Exception:
        if last_error:
            raise last_error
        raise


def search_text_file(path, pattern):
    text = read_text_file(path)
    rows = []
    for _, _, snippet in find_matches(text, pattern):
        rows.append(("본문", snippet))
    return rows


def search_pdf(path, pattern):
    if fitz is None:
        raise RuntimeError("PyMuPDF가 설치되지 않았습니다.")

    rows = []
    with fitz.open(path) as doc:
        for page_index, page in enumerate(doc):
            text = page.get_text("text")
            for _, _, snippet in find_matches(text, pattern):
                rows.append((f"{page_index + 1}페이지", snippet))
                if len(rows) >= MAX_HITS_PER_FILE:
                    return rows
    return rows


def search_xlsx(path, pattern):
    if openpyxl is None:
        raise RuntimeError("openpyxl이 설치되지 않았습니다.")

    rows = []
    workbook = openpyxl.load_workbook(
        path,
        read_only=True,
        data_only=False
    )

    try:
        for worksheet in workbook.worksheets:
            for row in worksheet.iter_rows():
                for cell in row:
                    value = cell.value
                    if value is None:
                        continue
                    text = str(value)
                    matches = find_matches(text, pattern)
                    for _, _, snippet in matches:
                        rows.append((f"{worksheet.title}!{cell.coordinate}", snippet))
                        if len(rows) >= MAX_HITS_PER_FILE:
                            return rows
    finally:
        workbook.close()

    return rows


def search_xls(path, pattern):
    if xlrd is None:
        raise RuntimeError("xlrd가 설치되지 않았습니다.")

    rows = []
    workbook = xlrd.open_workbook(path, on_demand=True)
    try:
        for sheet in workbook.sheets():
            for row_index in range(sheet.nrows):
                for col_index in range(sheet.ncols):
                    value = sheet.cell_value(row_index, col_index)
                    if value in (None, ""):
                        continue
                    text = str(value)
                    for _, _, snippet in find_matches(text, pattern):
                        cell_name = f"{xlrd.formula.colname(col_index)}{row_index + 1}"
                        rows.append((f"{sheet.name}!{cell_name}", snippet))
                        if len(rows) >= MAX_HITS_PER_FILE:
                            return rows
    finally:
        workbook.release_resources()

    return rows


def search_docx(path, pattern):
    if Document is None:
        raise RuntimeError("python-docx가 설치되지 않았습니다.")

    rows = []
    doc = Document(path)

    for index, paragraph in enumerate(doc.paragraphs, start=1):
        text = paragraph.text
        for _, _, snippet in find_matches(text, pattern):
            rows.append((f"문단 {index}", snippet))
            if len(rows) >= MAX_HITS_PER_FILE:
                return rows

    for table_index, table in enumerate(doc.tables, start=1):
        for row_index, row in enumerate(table.rows, start=1):
            for col_index, cell in enumerate(row.cells, start=1):
                text = cell.text
                for _, _, snippet in find_matches(text, pattern):
                    rows.append((f"표 {table_index} / {row_index}행 {col_index}열", snippet))
                    if len(rows) >= MAX_HITS_PER_FILE:
                        return rows

    return rows


def extract_shape_text(shape):
    """PowerPoint 도형에서 텍스트와 표 내용을 추출한다."""
    texts = []

    if hasattr(shape, "text") and shape.text:
        texts.append(shape.text)

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text:
                    texts.append(cell.text)

    if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):  # 그룹 도형
        for child in shape.shapes:
            texts.extend(extract_shape_text(child))

    return texts


def search_pptx(path, pattern):
    if Presentation is None:
        raise RuntimeError("python-pptx가 설치되지 않았습니다.")

    rows = []
    presentation = Presentation(path)

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            for text in extract_shape_text(shape):
                for _, _, snippet in find_matches(text, pattern):
                    rows.append((f"{slide_index}슬라이드", snippet))
                    if len(rows) >= MAX_HITS_PER_FILE:
                        return rows

    return rows


def search_hwpx(path, pattern):
    """
    HWPX는 ZIP 기반 XML 형식이다.
    XML에 포함된 텍스트를 추출해 검색한다.
    """
    rows = []

    with zipfile.ZipFile(path, "r") as archive:
        xml_names = [
            name for name in archive.namelist()
            if name.lower().endswith(".xml")
            and (
                name.startswith("Contents/")
                or name.startswith("Preview/")
            )
        ]

        for name in xml_names:
            raw = archive.read(name)
            try:
                root = ET.fromstring(raw)
                text = " ".join(part for part in root.itertext() if part)
            except ET.ParseError:
                text = raw.decode("utf-8", errors="ignore")

            for _, _, snippet in find_matches(text, pattern):
                location = Path(name).name
                rows.append((f"HWPX/{location}", snippet))
                if len(rows) >= MAX_HITS_PER_FILE:
                    return rows

    return rows


def search_file_contents(path, pattern):
    ext = path.suffix.lower()

    if ext in TEXT_EXTENSIONS:
        return search_text_file(path, pattern)
    if ext == ".pdf":
        return search_pdf(path, pattern)
    if ext in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
        return search_xlsx(path, pattern)
    if ext == ".xls":
        return search_xls(path, pattern)
    if ext == ".docx":
        return search_docx(path, pattern)
    if ext == ".pptx":
        return search_pptx(path, pattern)
    if ext == ".hwpx":
        return search_hwpx(path, pattern)

    return []


def build_pattern(keyword, use_regex, case_sensitive):
    flags = 0 if case_sensitive else re.IGNORECASE
    expression = keyword if use_regex else re.escape(keyword)
    return re.compile(expression, flags)


def parse_extension_filter(text):
    """
    "txt, .pdf *.docx" 같은 입력을 {'.txt', '.pdf', '.docx'} 형태로 정규화한다.
    비어 있으면 빈 집합(필터 없음, 전체 검색)을 반환한다.
    """
    if not text:
        return set()

    tokens = re.split(r"[,;\s]+", text.strip())
    extensions = set()
    for token in tokens:
        token = token.strip().lstrip("*")
        if not token:
            continue
        if not token.startswith("."):
            token = "." + token
        extensions.add(token.lower())
    return extensions


def open_path(path):
    """Windows에서는 기본 연결 프로그램으로 파일을 연다."""
    path = str(path)
    if sys.platform.startswith("win"):
        os.startfile(path)
    elif sys.platform == "darwin":
        subprocess.Popen(["open", path])
    else:
        subprocess.Popen(["xdg-open", path])


class ContentSearchApp(tk.Tk):
    def __init__(self):
        super().__init__()

        self.title(APP_TITLE)
        self.geometry("1280x760")
        self.minsize(980, 600)

        self.task_queue = queue.Queue()
        self.cancel_event = threading.Event()
        self.worker = None
        self.result_paths = {}

        self.folder_var = tk.StringVar()
        self.keyword_var = tk.StringVar()
        self.extension_var = tk.StringVar()
        self.subfolder_var = tk.BooleanVar(value=True)
        self.case_sensitive_var = tk.BooleanVar(value=False)
        self.regex_var = tk.BooleanVar(value=False)
        self.filename_var = tk.BooleanVar(value=True)
        self.contents_var = tk.BooleanVar(value=True)
        self.status_var = tk.StringVar(value="폴더와 키워드를 입력하세요.")

        self._build_ui()
        self.after(100, self._process_queue)

    def _build_ui(self):
        outer = ttk.Frame(self, padding=12)
        outer.pack(fill="both", expand=True)

        search_frame = ttk.LabelFrame(outer, text="검색 조건", padding=10)
        search_frame.pack(fill="x")

        ttk.Label(search_frame, text="검색 폴더").grid(row=0, column=0, sticky="w", padx=(0, 8), pady=5)
        folder_entry = ttk.Entry(search_frame, textvariable=self.folder_var)
        folder_entry.grid(row=0, column=1, sticky="ew", pady=5)
        ttk.Button(search_frame, text="폴더 선택", command=self._select_folder).grid(
            row=0, column=2, padx=(8, 0), pady=5
        )

        ttk.Label(search_frame, text="검색 키워드").grid(row=1, column=0, sticky="w", padx=(0, 8), pady=5)
        keyword_entry = ttk.Entry(search_frame, textvariable=self.keyword_var)
        keyword_entry.grid(row=1, column=1, sticky="ew", pady=5)
        keyword_entry.bind("<Return>", lambda event: self.start_search())

        ttk.Label(search_frame, text="확장자 필터").grid(row=2, column=0, sticky="w", padx=(0, 8), pady=5)
        extension_entry = ttk.Entry(search_frame, textvariable=self.extension_var)
        extension_entry.grid(row=2, column=1, sticky="ew", pady=5)
        extension_entry.bind("<Return>", lambda event: self.start_search())
        ttk.Label(search_frame, text="예: txt, pdf (비워두면 전체)", foreground="gray").grid(
            row=2, column=2, sticky="w", padx=(8, 0), pady=5
        )

        option_frame = ttk.Frame(search_frame)
        option_frame.grid(row=3, column=1, columnspan=2, sticky="w", pady=(4, 0))

        ttk.Checkbutton(option_frame, text="하위 폴더 포함", variable=self.subfolder_var).pack(side="left", padx=(0, 12))
        ttk.Checkbutton(option_frame, text="대소문자 구분", variable=self.case_sensitive_var).pack(side="left", padx=(0, 12))
        ttk.Checkbutton(option_frame, text="정규식 사용", variable=self.regex_var).pack(side="left", padx=(0, 12))
        ttk.Checkbutton(option_frame, text="파일명 검색", variable=self.filename_var).pack(side="left", padx=(0, 12))
        ttk.Checkbutton(option_frame, text="파일 내부 검색", variable=self.contents_var).pack(side="left")

        search_frame.columnconfigure(1, weight=1)

        button_frame = ttk.Frame(outer)
        button_frame.pack(fill="x", pady=10)

        self.search_button = ttk.Button(button_frame, text="검색 시작", command=self.start_search)
        self.search_button.pack(side="left")

        self.stop_button = ttk.Button(button_frame, text="검색 중지", command=self.stop_search, state="disabled")
        self.stop_button.pack(side="left", padx=(8, 0))

        ttk.Button(button_frame, text="결과 초기화", command=self.clear_results).pack(side="left", padx=(8, 0))
        ttk.Button(button_frame, text="CSV 저장", command=self.export_csv).pack(side="left", padx=(8, 0))

        self.progress = ttk.Progressbar(button_frame, mode="determinate")
        self.progress.pack(side="right", fill="x", expand=True, padx=(20, 0))

        result_frame = ttk.LabelFrame(outer, text="검색 결과", padding=8)
        result_frame.pack(fill="both", expand=True)

        columns = ("name", "type", "location", "snippet", "path")
        self.tree = ttk.Treeview(result_frame, columns=columns, show="headings", selectmode="browse")

        self.tree.heading("name", text="파일명")
        self.tree.heading("type", text="형식")
        self.tree.heading("location", text="일치 위치")
        self.tree.heading("snippet", text="일치 내용")
        self.tree.heading("path", text="전체 경로")

        self.tree.column("name", width=220, minwidth=140)
        self.tree.column("type", width=70, minwidth=60, anchor="center")
        self.tree.column("location", width=150, minwidth=100)
        self.tree.column("snippet", width=430, minwidth=220)
        self.tree.column("path", width=420, minwidth=220)

        y_scroll = ttk.Scrollbar(result_frame, orient="vertical", command=self.tree.yview)
        x_scroll = ttk.Scrollbar(result_frame, orient="horizontal", command=self.tree.xview)
        self.tree.configure(yscrollcommand=y_scroll.set, xscrollcommand=x_scroll.set)

        self.tree.grid(row=0, column=0, sticky="nsew")
        y_scroll.grid(row=0, column=1, sticky="ns")
        x_scroll.grid(row=1, column=0, sticky="ew")

        result_frame.rowconfigure(0, weight=1)
        result_frame.columnconfigure(0, weight=1)

        self.tree.bind("<Double-1>", lambda event: self.open_selected_file())
        self.tree.bind("<Return>", lambda event: self.open_selected_file())

        action_frame = ttk.Frame(outer)
        action_frame.pack(fill="x", pady=(8, 0))

        ttk.Button(action_frame, text="선택 파일 열기", command=self.open_selected_file).pack(side="left")
        ttk.Button(action_frame, text="선택 파일 폴더 열기", command=self.open_selected_folder).pack(side="left", padx=(8, 0))

        ttk.Label(action_frame, textvariable=self.status_var).pack(side="right")

    def _select_folder(self):
        folder = filedialog.askdirectory(title="검색할 폴더 선택")
        if folder:
            self.folder_var.set(folder)

    def clear_results(self):
        if self.worker and self.worker.is_alive():
            messagebox.showwarning("검색 중", "검색을 중지한 뒤 결과를 초기화하세요.")
            return

        for item in self.tree.get_children():
            self.tree.delete(item)
        self.result_paths.clear()
        self.status_var.set("결과를 초기화했습니다.")
        self.progress.stop()
        self.progress.configure(mode="determinate")
        self.progress["value"] = 0

    def start_search(self):
        folder = self.folder_var.get().strip()
        keyword = self.keyword_var.get()

        if not folder:
            messagebox.showwarning("입력 확인", "검색할 폴더를 선택하세요.")
            return

        if not Path(folder).is_dir():
            messagebox.showerror("폴더 오류", "선택한 폴더가 존재하지 않습니다.")
            return

        if not keyword:
            messagebox.showwarning("입력 확인", "검색 키워드를 입력하세요.")
            return

        if not self.filename_var.get() and not self.contents_var.get():
            messagebox.showwarning("검색 범위", "파일명 검색 또는 파일 내부 검색 중 하나 이상을 선택하세요.")
            return

        try:
            pattern = build_pattern(
                keyword,
                self.regex_var.get(),
                self.case_sensitive_var.get()
            )
        except re.error as exc:
            messagebox.showerror("정규식 오류", f"정규식이 올바르지 않습니다.\n\n{exc}")
            return

        extensions = parse_extension_filter(self.extension_var.get())

        self.clear_results()
        self.cancel_event.clear()
        self.search_button.configure(state="disabled")
        self.stop_button.configure(state="normal")
        self.status_var.set("검색을 시작합니다...")
        self.progress.configure(mode="indeterminate")
        self.progress.start(10)

        options = {
            "folder": Path(folder),
            "pattern": pattern,
            "include_subfolders": self.subfolder_var.get(),
            "search_filename": self.filename_var.get(),
            "search_contents": self.contents_var.get(),
            "extensions": extensions,
        }

        self.worker = threading.Thread(
            target=self._search_worker,
            args=(options,),
            daemon=True
        )
        self.worker.start()

    def stop_search(self):
        self.cancel_event.set()
        self.status_var.set("검색 중지 요청을 처리하는 중입니다...")

    def _iter_files(self, folder, include_subfolders, extensions=None):
        """폴더를 순회하며 파일을 찾는 즉시 하나씩 내보낸다(전체 목록을 먼저 모으지 않음)."""
        iterator = folder.rglob("*") if include_subfolders else folder.glob("*")
        for path in iterator:
            if self.cancel_event.is_set():
                return
            try:
                if not path.is_file():
                    continue
                if extensions and path.suffix.lower() not in extensions:
                    continue
                yield path
            except OSError:
                continue

    def _search_worker(self, options):
        folder = options["folder"]
        pattern = options["pattern"]
        include_subfolders = options["include_subfolders"]
        search_filename = options["search_filename"]
        search_contents = options["search_contents"]
        extensions = options["extensions"]

        matched_files = set()
        total_hits = 0
        processed = 0
        errors = []

        try:
            for path in self._iter_files(folder, include_subfolders, extensions):
                if self.cancel_event.is_set():
                    self.task_queue.put(("cancelled", processed, total_hits, len(matched_files), errors))
                    return

                processed += 1
                self.task_queue.put(("progress", processed, str(path)))

                try:
                    # 파일명 검색
                    if search_filename:
                        filename_matches = find_matches(path.name, pattern)
                        for _, _, snippet in filename_matches:
                            total_hits += 1
                            matched_files.add(str(path))
                            self.task_queue.put((
                                "result",
                                path.name,
                                path.suffix.lower() or "(없음)",
                                "파일명",
                                snippet,
                                str(path)
                            ))

                    # 파일 내부 검색
                    if search_contents and path.suffix.lower() in SUPPORTED_EXTENSIONS:
                        content_rows = search_file_contents(path, pattern)
                        for location, snippet in content_rows:
                            total_hits += 1
                            matched_files.add(str(path))
                            self.task_queue.put((
                                "result",
                                path.name,
                                path.suffix.lower() or "(없음)",
                                location,
                                snippet,
                                str(path)
                            ))

                except PermissionError:
                    errors.append(f"권한 없음: {path}")
                except Exception as exc:
                    errors.append(f"{path} | {exc}")

            if self.cancel_event.is_set():
                self.task_queue.put(("cancelled", processed, total_hits, len(matched_files), errors))
            else:
                self.task_queue.put(("done", processed, total_hits, len(matched_files), errors))

        except Exception as exc:
            self.task_queue.put(("fatal", str(exc), traceback.format_exc()))

    def _process_queue(self):
        try:
            while True:
                event = self.task_queue.get_nowait()
                kind = event[0]

                if kind == "progress":
                    index, path = event[1], event[2]
                    self.status_var.set(f"{index:,}개 확인 중: {Path(path).name}")

                elif kind == "result":
                    _, name, ext, location, snippet, path = event
                    iid = self.tree.insert(
                        "",
                        "end",
                        values=(name, ext, location, snippet, path)
                    )
                    self.result_paths[iid] = path

                elif kind == "done":
                    processed, hit_count, file_count, errors = event[1], event[2], event[3], event[4]
                    self.search_button.configure(state="normal")
                    self.stop_button.configure(state="disabled")
                    self.progress.stop()
                    self.progress.configure(mode="determinate")
                    self.progress["value"] = self.progress["maximum"]
                    self.status_var.set(
                        f"완료: {processed:,}개 파일 검사 / {file_count:,}개 파일 일치 / {hit_count:,}건 발견"
                    )
                    if errors:
                        self._show_error_summary(errors)

                elif kind == "cancelled":
                    processed, hit_count, file_count, errors = event[1], event[2], event[3], event[4]
                    self.search_button.configure(state="normal")
                    self.stop_button.configure(state="disabled")
                    self.progress.stop()
                    self.progress.configure(mode="determinate")
                    self.progress["value"] = 0
                    self.status_var.set(
                        f"중지됨: {processed:,}개 검사 / {file_count:,}개 파일 일치 / {hit_count:,}건 발견"
                    )
                    if errors:
                        self._show_error_summary(errors)

                elif kind == "fatal":
                    message, detail = event[1], event[2]
                    self.search_button.configure(state="normal")
                    self.stop_button.configure(state="disabled")
                    self.progress.stop()
                    self.progress.configure(mode="determinate")
                    self.progress["value"] = 0
                    self.status_var.set("오류로 검색이 중단되었습니다.")
                    messagebox.showerror("검색 오류", f"{message}\n\n{detail}")

        except queue.Empty:
            pass

        self.after(100, self._process_queue)

    def _show_error_summary(self, errors):
        preview = "\n".join(errors[:10])
        more = ""
        if len(errors) > 10:
            more = f"\n\n외 {len(errors) - 10:,}건"
        messagebox.showwarning(
            "일부 파일 검색 실패",
            "일부 파일은 손상, 암호화, 권한 또는 형식 문제로 검색하지 못했습니다.\n\n"
            + preview
            + more
        )

    def _selected_path(self):
        selection = self.tree.selection()
        if not selection:
            messagebox.showinfo("선택 필요", "검색 결과에서 파일을 선택하세요.")
            return None

        item_id = selection[0]
        return self.result_paths.get(item_id)

    def open_selected_file(self):
        path = self._selected_path()
        if not path:
            return

        try:
            open_path(path)
        except Exception as exc:
            messagebox.showerror("파일 열기 실패", str(exc))

    def open_selected_folder(self):
        path = self._selected_path()
        if not path:
            return

        try:
            if sys.platform.startswith("win"):
                subprocess.Popen(["explorer", "/select,", os.path.normpath(path)])
            else:
                open_path(Path(path).parent)
        except Exception as exc:
            messagebox.showerror("폴더 열기 실패", str(exc))

    def export_csv(self):
        items = self.tree.get_children()
        if not items:
            messagebox.showinfo("저장할 결과 없음", "먼저 검색을 실행하세요.")
            return

        save_path = filedialog.asksaveasfilename(
            title="검색 결과 CSV 저장",
            defaultextension=".csv",
            filetypes=[("CSV 파일", "*.csv")]
        )
        if not save_path:
            return

        try:
            with open(save_path, "w", newline="", encoding="utf-8-sig") as f:
                writer = csv.writer(f)
                writer.writerow(["파일명", "형식", "일치 위치", "일치 내용", "전체 경로"])
                for item in items:
                    writer.writerow(self.tree.item(item, "values"))

            messagebox.showinfo("저장 완료", f"CSV 파일을 저장했습니다.\n\n{save_path}")
        except Exception as exc:
            messagebox.showerror("저장 실패", str(exc))


if __name__ == "__main__":
    app = ContentSearchApp()
    app.mainloop()
