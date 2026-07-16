"""파일 형식별 텍스트 추출 및 검색. 선택적 라이브러리가 없으면 해당 형식은 건너뛴다."""

from pathlib import Path
from xml.etree import ElementTree as ET
import zipfile

from .text_utils import MAX_HITS_PER_FILE, find_matches

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import xlrd
except ImportError:
    xlrd = None


TEXT_EXTENSIONS = {
    ".txt", ".csv", ".tsv", ".log", ".ini", ".cfg", ".conf",
    ".json", ".xml", ".yaml", ".yml", ".md", ".py", ".java",
    ".c", ".cpp", ".h", ".hpp", ".js", ".ts", ".css", ".html",
    ".htm", ".sql", ".bat", ".cmd", ".ps1"
}

SUPPORTED_EXTENSIONS = (
    TEXT_EXTENSIONS
    | {".pdf", ".xlsx", ".xlsm", ".xltx", ".xltm", ".xls",
       ".docx", ".pptx", ".hwpx"}
)


def read_text_file(path):
    """텍스트 파일을 여러 인코딩으로 시도해 읽는다."""
    encodings = ("utf-8-sig", "utf-8", "cp949", "euc-kr", "utf-16")
    last_error = None

    for encoding in encodings:
        try:
            with open(path, "r", encoding=encoding, errors="strict") as f:
                return f.read()
        except UnicodeError as exc:
            last_error = exc

    # 마지막 수단: 읽을 수 없는 문자를 대체
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
    except Exception:
        if last_error:
            raise last_error
        raise


# ----------------------------------------------------------------------
# 각 형식에서 (위치, 텍스트) 단위를 순서대로 꺼내는 제너레이터.
# search_xxx()는 여기서 찾은 텍스트에 검색 패턴을 적용해 결과 행을 만들고,
# file_contains_excluded_text()는 같은 제너레이터로 제외 키워드가 파일
# 어딘가에 있는지만 빠르게 확인한다(두 기능이 파일 파싱 로직을 공유한다).
# ----------------------------------------------------------------------

def _iter_text_file_units(path):
    yield ("본문", read_text_file(path))


def _iter_pdf_units(path):
    if fitz is None:
        raise RuntimeError("PyMuPDF가 설치되지 않았습니다.")
    with fitz.open(path) as doc:
        for page_index, page in enumerate(doc):
            yield (f"{page_index + 1}페이지", page.get_text("text"))


def _iter_xlsx_units(path):
    if openpyxl is None:
        raise RuntimeError("openpyxl이 설치되지 않았습니다.")
    workbook = openpyxl.load_workbook(path, read_only=True, data_only=False)
    try:
        for worksheet in workbook.worksheets:
            for row in worksheet.iter_rows():
                for cell in row:
                    value = cell.value
                    if value is None:
                        continue
                    yield (f"{worksheet.title}!{cell.coordinate}", str(value))
    finally:
        workbook.close()


def _iter_xls_units(path):
    if xlrd is None:
        raise RuntimeError("xlrd가 설치되지 않았습니다.")
    workbook = xlrd.open_workbook(path, on_demand=True)
    try:
        for sheet in workbook.sheets():
            for row_index in range(sheet.nrows):
                for col_index in range(sheet.ncols):
                    value = sheet.cell_value(row_index, col_index)
                    if value in (None, ""):
                        continue
                    cell_name = f"{xlrd.formula.colname(col_index)}{row_index + 1}"
                    yield (f"{sheet.name}!{cell_name}", str(value))
    finally:
        workbook.release_resources()


def _iter_docx_units(path):
    if Document is None:
        raise RuntimeError("python-docx가 설치되지 않았습니다.")
    doc = Document(path)

    for index, paragraph in enumerate(doc.paragraphs, start=1):
        yield (f"문단 {index}", paragraph.text)

    for table_index, table in enumerate(doc.tables, start=1):
        for row_index, row in enumerate(table.rows, start=1):
            for col_index, cell in enumerate(row.cells, start=1):
                yield (f"표 {table_index} / {row_index}행 {col_index}열", cell.text)


def extract_shape_text(shape):
    """PowerPoint 도형에서 텍스트와 표 내용을 추출한다."""
    texts = []

    if hasattr(shape, "text") and shape.text:
        texts.append(shape.text)

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text:
                    texts.append(cell.text)

    if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):  # 그룹 도형
        for child in shape.shapes:
            texts.extend(extract_shape_text(child))

    return texts


def _iter_pptx_units(path):
    if Presentation is None:
        raise RuntimeError("python-pptx가 설치되지 않았습니다.")
    presentation = Presentation(path)

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            for text in extract_shape_text(shape):
                yield (f"{slide_index}슬라이드", text)


def _iter_hwpx_units(path):
    """HWPX는 ZIP 기반 XML 형식이다. XML에 포함된 텍스트를 추출한다."""
    with zipfile.ZipFile(path, "r") as archive:
        xml_names = [
            name for name in archive.namelist()
            if name.lower().endswith(".xml")
            and (
                name.startswith("Contents/")
                or name.startswith("Preview/")
            )
        ]

        for name in xml_names:
            raw = archive.read(name)
            try:
                root = ET.fromstring(raw)
                text = " ".join(part for part in root.itertext() if part)
            except ET.ParseError:
                text = raw.decode("utf-8", errors="ignore")

            yield (f"HWPX/{Path(name).name}", text)


_UNIT_ITERATORS = {ext: _iter_text_file_units for ext in TEXT_EXTENSIONS}
_UNIT_ITERATORS[".pdf"] = _iter_pdf_units
for _ext in (".xlsx", ".xlsm", ".xltx", ".xltm"):
    _UNIT_ITERATORS[_ext] = _iter_xlsx_units
_UNIT_ITERATORS[".xls"] = _iter_xls_units
_UNIT_ITERATORS[".docx"] = _iter_docx_units
_UNIT_ITERATORS[".pptx"] = _iter_pptx_units
_UNIT_ITERATORS[".hwpx"] = _iter_hwpx_units


def _search_units(iter_units, path, pattern):
    rows = []
    for location, text in iter_units(path):
        for _, _, keyword, snippet in find_matches(text, pattern):
            rows.append((location, keyword, snippet))
            if len(rows) >= MAX_HITS_PER_FILE:
                return rows
    return rows


def search_text_file(path, pattern):
    return _search_units(_iter_text_file_units, path, pattern)


def search_pdf(path, pattern):
    return _search_units(_iter_pdf_units, path, pattern)


def search_xlsx(path, pattern):
    return _search_units(_iter_xlsx_units, path, pattern)


def search_xls(path, pattern):
    return _search_units(_iter_xls_units, path, pattern)


def search_docx(path, pattern):
    return _search_units(_iter_docx_units, path, pattern)


def search_pptx(path, pattern):
    return _search_units(_iter_pptx_units, path, pattern)


def search_hwpx(path, pattern):
    return _search_units(_iter_hwpx_units, path, pattern)


def search_file_contents(path, pattern):
    ext = path.suffix.lower()

    if ext in TEXT_EXTENSIONS:
        return search_text_file(path, pattern)
    if ext == ".pdf":
        return search_pdf(path, pattern)
    if ext in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
        return search_xlsx(path, pattern)
    if ext == ".xls":
        return search_xls(path, pattern)
    if ext == ".docx":
        return search_docx(path, pattern)
    if ext == ".pptx":
        return search_pptx(path, pattern)
    if ext == ".hwpx":
        return search_hwpx(path, pattern)

    return []


def file_contains_excluded_text(path, exclude_pattern):
    """
    이 파일의 지원되는 내용 어딘가(페이지/셀/문단/슬라이드/XML 등)에
    exclude_pattern이 하나라도 있으면 True를 반환한다. 파일명은 검사하지
    않는다(파일명은 호출하는 쪽에서 별도로 확인한다).
    """
    iterator = _UNIT_ITERATORS.get(path.suffix.lower())
    if iterator is None:
        return False

    for _location, text in iterator(path):
        if exclude_pattern.search(text):
            return True
    return False
